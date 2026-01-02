import os
import requests
from pypdf import PdfReader
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# ===============================
# CONFIG
# ===============================
BASE_PATH = "data"
CHUNK_SIZE = 400
TOP_K = 3
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

# ===============================
# STORAGE
# ===============================
documents = []
sources = []
profile_chunks = []

# ===============================
# HELPER: CHUNK TEXT
# ===============================
def add_chunks(text, source):
    for i in range(0, len(text), CHUNK_SIZE):
        chunk = text[i:i + CHUNK_SIZE]
        if chunk.strip():
            documents.append(chunk)
            sources.append(source)

# ===============================
# LOAD PDFs & PPTs
# ===============================
for subject in os.listdir(BASE_PATH):
    subject_path = os.path.join(BASE_PATH, subject)
    if not os.path.isdir(subject_path):
        continue

    for file in os.listdir(subject_path):
        path = os.path.join(subject_path, file)

        # ---------- PDFs ----------
        if file.lower().endswith(".pdf"):
            reader = PdfReader(path)
            for page_no, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    add_chunks(text, f"{file} (PDF page {page_no + 1})")

        # ---------- PPTs ----------
        elif file.lower().endswith(".pptx"):
            prs = Presentation(path)
            for slide_no, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                if slide_text.strip():
                    add_chunks(slide_text, f"{file} (PPT slide {slide_no + 1})")

# ===============================
# LOAD PROFILE
# ===============================
if os.path.exists("profile.txt"):
    with open("profile.txt", "r", encoding="utf-8") as f:
        profile_text = f.read()
        for i in range(0, len(profile_text), CHUNK_SIZE):
            chunk = profile_text[i:i + CHUNK_SIZE]
            if chunk.strip():
                profile_chunks.append(chunk)
                documents.append(chunk)
                sources.append("User Profile")

if len(documents) == 0:
    raise ValueError("❌ No PDFs, PPTs, or profile.txt found.")

# ===============================
# VECTORIZE (TF-IDF)
# ===============================
vectorizer = TfidfVectorizer(stop_words="english")
doc_vectors = vectorizer.fit_transform(documents)

# ===============================
# GROQ LLaMA GENERATION (SAFE)
# ===============================
def groq_generate(context, query):
    if not GROQ_API_KEY:
        return "❌ GROQ_API_KEY not set. Please set it and restart VS Code."

    prompt = f"""
Use the context below to answer clearly and concisely.

Context:
{context}

Question:
{query}
"""

    response = requests.post(
        "https://api.groq.com/openai/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        },
        json={
            "model": "llama-3.1-8b-instant",
            "messages": [
                {"role": "system", "content": "You are an academic assistant."},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.2
        }
    )

    try:
        data = response.json()
    except Exception:
        return "⚠️ Failed to parse Groq response."

    if "choices" not in data:
        return "⚠️ Groq API error:\n\n" + str(data)

    return data["choices"][0]["message"]["content"]

# ===============================
# ASK FUNCTION
# ===============================
def ask(query):
    q = query.lower()

    # ---------- PROFILE PRIORITY ----------
    if any(
        keyword in q
        for keyword in [
            "who is",
            "who's",
            "waseeq",
            "siddique",
            "profile",
            "student",
            "fyp"
        ]
    ):
        if profile_chunks:
            context = "\n".join(profile_chunks)
            return groq_generate(context, query)

    # ---------- NORMAL RETRIEVAL ----------
    query_vec = vectorizer.transform([query])
    similarities = cosine_similarity(query_vec, doc_vectors)[0]
    top_indices = similarities.argsort()[-TOP_K:][::-1]

    context = "\n".join([documents[i] for i in top_indices])
    return groq_generate(context, query)